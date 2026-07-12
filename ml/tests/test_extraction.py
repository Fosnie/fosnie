"""Native office extraction (XLSX/PPTX) + OCR routing for images and scanned
PDFs. The OCR service call is mocked — no GLM-OCR
network call. Office/text formats must NEVER touch OCR; a text-bearing PDF uses
its embedded layer; only images and scanned (no-text) PDFs are sent to OCR."""

import asyncio

import pytest

from app import extract, ocr
from app.config import settings


def test_xlsx_native_extraction(tmp_path):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Fees"
    ws.append(["Item", "Amount"])
    ws.append(["Consultancy", 30])
    p = tmp_path / "f.xlsx"
    wb.save(str(p))

    text = extract.extract(str(p))
    assert "Fees" in text and "Consultancy" in text and "30" in text


def test_xlsx_prepends_row_count_header(tmp_path):
    # A counting question must be answerable from metadata, not by tallying rows.
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "200iq"
    ws.append(["Title", "Artist", "Plays"])  # header
    for i in range(812):  # 812 data rows
        ws.append([f"Track {i}", f"Artist {i}", i])
    p = tmp_path / "playlist.xlsx"
    wb.save(str(p))

    first_line = extract.extract(str(p)).splitlines()[0]
    assert '[Sheet "200iq"' in first_line
    assert "813 rows" in first_line  # 1 header + 812 data
    assert "812 data rows" in first_line
    assert "3 columns" in first_line


def test_csv_prepends_row_count_header(tmp_path):
    p = tmp_path / "playlist.csv"
    p.write_text("Title,Artist\n" + "\n".join(f"Track {i},Artist {i}" for i in range(812)), encoding="utf-8")
    out = extract.extract(str(p))
    first_line = out.splitlines()[0]
    assert "[CSV: 813 rows" in first_line and "812 data rows" in first_line and "2 columns" in first_line
    assert "Track 0" in out  # original body retained


def test_pptx_native_extraction(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Quarterly terms"
    p = tmp_path / "d.pptx"
    prs.save(str(p))

    text = extract.extract(str(p))
    assert "Quarterly terms" in text


def test_image_routes_to_ocr(tmp_path, monkeypatch):
    seen = {}

    async def fake_ocr(data, mime):
        seen["mime"] = mime
        return "TEXT FROM IMAGE"

    monkeypatch.setattr(ocr, "ocr_bytes", fake_ocr)
    img = tmp_path / "scan.png"
    img.write_bytes(b"\x89PNG fake")
    pages = asyncio.run(extract.extract_pages_ocr(str(img)))
    assert pages == [(1, "TEXT FROM IMAGE")]
    assert seen["mime"] == "image/png"


def test_scanned_pdf_routes_to_ocr(tmp_path, monkeypatch):
    # No text layer → OCR the whole PDF (the service rasterises).
    monkeypatch.setattr(extract, "extract_pages", lambda path, mime=None: [(1, "")])

    async def fake_ocr(data, mime):
        assert mime == "application/pdf"
        return "OCR TEXT"

    monkeypatch.setattr(ocr, "ocr_bytes", fake_ocr)
    pdf = tmp_path / "scan.pdf"
    pdf.write_bytes(b"%PDF-1.4 fake")
    pages = asyncio.run(extract.extract_pages_ocr(str(pdf)))
    assert pages == [(1, "OCR TEXT")]


def test_text_pdf_uses_layer_not_ocr(tmp_path, monkeypatch):
    monkeypatch.setattr(extract, "extract_pages", lambda path, mime=None: [(1, "real text layer here")])

    async def boom(data, mime):
        raise AssertionError("OCR must not run for a text-bearing PDF")

    monkeypatch.setattr(ocr, "ocr_bytes", boom)
    pdf = tmp_path / "t.pdf"
    pdf.write_bytes(b"%PDF fake")
    pages = asyncio.run(extract.extract_pages_ocr(str(pdf)))
    assert pages == [(1, "real text layer here")]


def test_office_path_never_ocrs(tmp_path, monkeypatch):
    monkeypatch.setattr(extract, "extract_pages", lambda path, mime=None: [(1, "doc text")])

    async def boom(data, mime):
        raise AssertionError("OCR must not run for office/text formats")

    monkeypatch.setattr(ocr, "ocr_bytes", boom)
    pages = asyncio.run(extract.extract_pages_ocr("whatever.docx"))
    assert pages == [(1, "doc text")]


def test_ocr_disabled_raises(monkeypatch):
    monkeypatch.setattr(settings, "ocr_enabled", False)
    with pytest.raises(ocr.OcrUnavailable):
        asyncio.run(ocr.ocr_bytes(b"data", "image/png"))


def test_ocr_empty_input_raises(monkeypatch):
    monkeypatch.setattr(settings, "ocr_enabled", True)
    with pytest.raises(ocr.OcrUnavailable):
        asyncio.run(ocr.ocr_bytes(b"", "image/png"))
