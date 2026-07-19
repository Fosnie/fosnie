"""Presentation artefacts: the slide-spec parser and its Markdown fallback, every
layout archetype, the overflow budgets (font ladder + continuation slides), the
native chart contract, speaker notes, and the deterministic validators."""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Pt

from app import generate, pptx_engine, validators

REPO_ROOT = Path(__file__).resolve().parents[2]


def _build(spec, tmp_path, name="deck") -> Presentation:
    content = spec if isinstance(spec, str) else json.dumps(spec)
    out = str(tmp_path / f"{name}.pptx")
    res = generate.generate_artefact("pptx", "Deck", content, out)
    assert res["mime"].endswith("presentationml.presentation")
    return Presentation(res["path"])


def _texts(slide) -> list[str]:
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            out.extend(p.text for p in shape.text_frame.paragraphs if p.text.strip())
    return out


def _runs(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                yield from para.runs


# --- parsing ------------------------------------------------------------------


def test_parse_accepts_the_three_json_shapes():
    slide = {"layout": "bullets", "title": "T", "bullets": ["a"]}
    assert len(pptx_engine._parse(json.dumps({"slides": [slide]}), "X")["slides"]) == 1
    assert len(pptx_engine._parse(json.dumps([slide, slide]), "X")["slides"]) == 2
    assert len(pptx_engine._parse(json.dumps(slide), "X")["slides"]) == 1  # a lone slide


def test_parse_reads_markdown_for_the_drafter_path():
    md = "# Board update\n\nWhere we stand.\n\n## Findings\n\n- One failed\n- Two passed\n"
    slides = pptx_engine._parse(md, "Fallback")["slides"]
    assert slides[0]["layout"] == "title"
    assert slides[0]["title"] == "Board update"
    assert slides[0]["subtitle"] == "Where we stand."   # the lead paragraph
    assert slides[1]["title"] == "Findings"
    assert slides[1]["bullets"] == ["One failed", "Two passed"]


def test_parse_strips_markdown_emphasis_in_the_fallback():
    slides = pptx_engine._parse("## Risk\n\n- **Overdue** reviews\n", "X")["slides"]
    assert slides[-1]["bullets"] == ["Overdue reviews"]


def test_parse_rejects_unusable_content():
    with pytest.raises(ValueError):
        pptx_engine._parse("", "X")
    with pytest.raises(ValueError):
        pptx_engine._parse('{"slides": []}', "X")
    with pytest.raises(ValueError):
        pptx_engine._parse('{"slides": [{"layout": "carousel"}]}', "X")  # unknown layout


def test_generate_rejects_garbage(tmp_path):
    with pytest.raises(ValueError):
        generate.generate_artefact("pptx", "x", "   ", str(tmp_path / "x.pptx"))


# --- layouts ------------------------------------------------------------------


def test_every_layout_builds_and_reopens(tmp_path):
    spec = {"slides": [
        {"layout": "title", "title": "Q2 Compliance Review", "subtitle": "Findings"},
        {"layout": "section", "title": "What we found"},
        {"layout": "bullets", "title": "Three controls failed",
         "bullets": ["Access reviews missed", {"text": "Retention lapsed", "sub": ["4,200 records"]}]},
        {"layout": "two_column", "title": "Before and after",
         "left": {"heading": "Now", "bullets": ["Manual"]},
         "right": {"heading": "Target", "bullets": ["Automated"]}},
        {"layout": "stat", "title": "In numbers",
         "stats": [{"value": "£1.2m", "label": "exposed"}, {"value": "38%", "label": "overdue"}]},
        {"layout": "table", "title": "Control status", "columns": ["Control", "State"],
         "rows": [["Access", "Failed"], ["Payments", "Passed"]], "caption": "Sample of four"},
        {"layout": "chart", "title": "Overdue doubled",
         "chart": {"type": "column", "categories": ["Jan", "Feb"],
                   "series": [{"name": "Overdue", "values": [9, 19]}]}},
        {"layout": "quote", "text": "The control did not operate.", "attribution": "Internal audit"},
        {"layout": "closing", "title": "Next steps", "bullets": ["Dual sign-off by September"]},
    ]}
    prs = _build(spec, tmp_path)
    assert len(prs.slides) == 9  # one slide per archetype, all reopened by python-pptx


def test_deck_is_widescreen_and_titles_are_live_text(tmp_path):
    prs = _build({"slides": [{"layout": "title", "title": "Board pack"}]}, tmp_path)
    assert round(prs.slide_width / 914400, 3) == 13.333  # 16:9
    assert round(prs.slide_height / 914400, 1) == 7.5
    assert "Board pack" in _texts(prs.slides[0])         # editable text, not an image


def test_section_slides_are_numbered_by_the_engine(tmp_path):
    spec = {"slides": [
        {"layout": "section", "title": "First"},
        {"layout": "bullets", "title": "Body", "bullets": ["x"]},
        {"layout": "section", "title": "Second"},
    ]}
    prs = _build(spec, tmp_path)
    assert "01" in _texts(prs.slides[0])
    assert "02" in _texts(prs.slides[2])


def test_table_and_chart_are_native_parts(tmp_path):
    spec = {"slides": [
        {"layout": "table", "title": "T", "columns": ["A", "B"], "rows": [["1", "2"]]},
        {"layout": "chart", "title": "C",
         "chart": {"type": "line", "categories": ["Jan"], "series": [{"name": "S", "values": [1]}]}},
    ]}
    prs = _build(spec, tmp_path)
    assert any(s.has_table for s in prs.slides[0].shapes)   # editable table, not a picture
    assert any(s.has_chart for s in prs.slides[1].shapes)   # editable chart data


def test_notes_land_in_the_notes_slide(tmp_path):
    spec = {"slides": [{"layout": "bullets", "title": "T", "bullets": ["x"],
                        "notes": "Say this out loud."}]}
    prs = _build(spec, tmp_path)
    assert prs.slides[0].notes_slide.notes_text_frame.text == "Say this out loud."


# --- budgets, ladders, continuation -------------------------------------------


def test_long_heading_steps_the_font_down(tmp_path):
    short = {"slides": [{"layout": "bullets", "title": "Short title", "bullets": ["x"]}]}
    long = {"slides": [{"layout": "bullets", "title": "W" * 130, "bullets": ["x"]}]}
    a = max(r.font.size for r in _runs(_build(short, tmp_path, "a").slides[0]) if r.font.size)
    b = max(r.font.size for r in _runs(_build(long, tmp_path, "b").slides[0]) if r.font.size)
    assert b < a                        # the ladder shrank the heading
    assert b >= Pt(22)                  # but never below the floor


def test_many_bullets_split_onto_a_continuation_slide(tmp_path):
    spec = {"slides": [{"layout": "bullets", "title": "Findings",
                        "bullets": [f"item {i}" for i in range(9)],
                        "notes": "only on the first part"}]}
    prs = _build(spec, tmp_path)
    assert len(prs.slides) == 2
    assert any("continued" in t for t in _texts(prs.slides[1]))
    assert prs.slides[0].has_notes_slide                      # notes stay with part one
    assert "item 8" in " ".join(_texts(prs.slides[1]))        # nothing is dropped


def test_six_bullets_stay_on_one_slide(tmp_path):
    spec = {"slides": [{"layout": "bullets", "title": "T",
                        "bullets": [f"item {i}" for i in range(6)]}]}
    assert len(_build(spec, tmp_path).slides) == 1  # the ladder handles it


def test_long_tables_split_and_repeat_the_header(tmp_path):
    spec = {"slides": [{"layout": "table", "title": "Log", "columns": ["Ref", "State"],
                        "rows": [[f"R{i}", "open"] for i in range(11)]}]}
    prs = _build(spec, tmp_path)
    assert len(prs.slides) == 2
    second = next(s.table for s in prs.slides[1].shapes if s.has_table)
    assert second.cell(0, 0).text == "Ref"                    # header repeated
    assert len(second.rows) == 4                              # 3 spilled rows + header


# --- charts -------------------------------------------------------------------


@pytest.mark.parametrize("ctype", ["bar", "column", "line", "pie", "doughnut"])
def test_each_chart_type_builds(ctype, tmp_path):
    spec = {"slides": [{"layout": "chart", "title": "C",
                        "chart": {"type": ctype, "categories": ["A", "B"],
                                  "series": [{"name": "S", "values": [1, 2]}]}}]}
    prs = _build(spec, tmp_path, ctype)
    assert any(s.has_chart for s in prs.slides[0].shapes)


def test_pie_rejects_more_than_one_series(tmp_path):
    spec = {"slides": [{"layout": "chart", "title": "C",
                        "chart": {"type": "pie", "categories": ["A"],
                                  "series": [{"name": "S1", "values": [1]},
                                             {"name": "S2", "values": [2]}]}}]}
    with pytest.raises(ValueError):
        _build(spec, tmp_path)


def test_chart_rejects_overlong_and_unusable_data(tmp_path):
    over = {"slides": [{"layout": "chart", "title": "C",
                        "chart": {"type": "column", "categories": [str(i) for i in range(9)],
                                  "series": [{"name": "S", "values": list(range(9))}]}}]}
    with pytest.raises(ValueError):
        _build(over, tmp_path, "over")

    words = {"slides": [{"layout": "chart", "title": "C",
                         "chart": {"type": "column", "categories": ["A"],
                                   "series": [{"name": "S", "values": ["£1.2m"]}]}}]}
    with pytest.raises(ValueError):
        _build(words, tmp_path, "words")

    unknown = {"slides": [{"layout": "chart", "title": "C",
                           "chart": {"type": "radar", "categories": ["A"],
                                     "series": [{"name": "S", "values": [1]}]}}]}
    with pytest.raises(ValueError):
        _build(unknown, tmp_path, "unknown")


# --- branding -----------------------------------------------------------------


def test_brand_colour_config_is_applied_and_typos_are_ignored(tmp_path, monkeypatch):
    from app.config import settings

    monkeypatch.setattr(settings, "pptx_colour_accent", "#AA0055", raising=False)
    prs = _build({"slides": [{"layout": "stat", "title": "T",
                              "stats": [{"value": "42", "label": "x"}]}]}, tmp_path, "brand")
    colours = {str(r.font.color.rgb) for r in _runs(prs.slides[0]) if r.font.color.type is not None}
    assert "AA0055" in colours

    monkeypatch.setattr(settings, "pptx_colour_accent", "not-a-colour", raising=False)
    prs = _build({"slides": [{"layout": "stat", "title": "T",
                              "stats": [{"value": "42", "label": "x"}]}]}, tmp_path, "typo")
    colours = {str(r.font.color.rgb) for r in _runs(prs.slides[0]) if r.font.color.type is not None}
    assert pptx_engine.ACCENT in colours   # fell back to the default rather than failing


# --- validators ---------------------------------------------------------------


def test_validate_pptx_passes_on_real_output(tmp_path):
    # Every layout, so the schema check covers the hand-built parts (bullet paragraph
    # formatting, tables, charts) and not just a bare title slide.
    spec = {"slides": [
        {"layout": "title", "title": "T", "subtitle": "S"},
        {"layout": "section", "title": "One"},
        {"layout": "bullets", "title": "B", "bullets": ["a", {"text": "b", "sub": ["c"]}]},
        {"layout": "two_column", "title": "C", "left": {"heading": "L", "bullets": ["x"]},
         "right": {"bullets": ["y"]}},
        {"layout": "stat", "title": "S", "stats": [{"value": "42", "label": "n"}]},
        {"layout": "table", "title": "T", "columns": ["A"], "rows": [["1"]], "caption": "cap"},
        {"layout": "chart", "title": "C", "caption": "cap",
         "chart": {"type": "pie", "categories": ["A", "B"],
                   "series": [{"name": "S", "values": [1, 2]}]}},
        {"layout": "quote", "text": "q", "attribution": "a"},
        {"layout": "closing", "title": "N", "bullets": ["z"]},
    ]}
    out = str(tmp_path / "ok.pptx")
    generate.generate_artefact("pptx", "Deck", json.dumps(spec), out)
    validators.validate_pptx(out)          # raises on failure
    validators.validate_pptx_xsd(out)      # schema-clean, not merely openable


def test_bullets_are_real_list_formatting(tmp_path):
    # The glyph belongs to the paragraph, not the text: wrapped lines hang correctly
    # and the recipient gets a genuine editable list rather than typed-in characters.
    from pptx.oxml.ns import qn

    prs = _build({"slides": [{"layout": "bullets", "title": "T",
                              "bullets": [{"text": "one", "sub": ["deep"]}]}]}, tmp_path)
    texts = _texts(prs.slides[0])
    assert "one" in texts and "deep" in texts       # no "• " baked into the string
    assert not any("•" in t or "–" in t for t in texts)
    bulleted = [
        para for shape in prs.slides[0].shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        if para._p.find(qn("a:pPr")) is not None
        and para._p.find(qn("a:pPr")).find(qn("a:buChar")) is not None
    ]
    assert len(bulleted) == 2                       # the bullet and its sub-bullet


def test_validate_pptx_rejects_garbage(tmp_path):
    p = tmp_path / "bad.pptx"
    p.write_bytes(b"not a deck")
    with pytest.raises(validators.ValidationError):
        validators.validate_pptx(str(p))


def test_validate_pptx_rejects_a_deck_without_slides(tmp_path):
    src = tmp_path / "src.pptx"
    generate.generate_artefact("pptx", "Deck", '{"slides":[{"layout":"title","title":"T"}]}', str(src))
    stripped = tmp_path / "empty.pptx"
    with zipfile.ZipFile(src) as z_in, zipfile.ZipFile(stripped, "w") as z_out:
        for item in z_in.infolist():
            if not item.filename.startswith("ppt/slides/slide"):
                z_out.writestr(item, z_in.read(item.filename))
    with pytest.raises(validators.ValidationError):
        validators.validate_pptx_xsd(str(stripped))


def test_pml_xsd_is_vendored():
    schemas = REPO_ROOT / "ml" / "app" / "assets" / "ooxml-schemas" / "ISO-IEC29500-4_2016"
    assert (schemas / "pml.xsd").is_file()
    assert (schemas / "dml-chart.xsd").is_file()   # charts are native OOXML parts
