# Copyright 2026 Private AI Ltd (SC881079)
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Text extraction. Native parsers for the office
formats — plain text, DOCX (shared tracked-change flattener), PDF text (pypdf),
XLSX (openpyxl), PPTX (python-pptx) — plus a vision-OCR path for scanned PDFs and
images, wired in via `extract_pages_ocr` (used by the async ingest pipeline).
The OCR engine (GLM-OCR over an OpenAI-compatible vision endpoint) is swappable
behind `ocr.py`."""

import logging
from pathlib import Path

import pypdf

from . import tracked_changes
from .config import settings
from .rag_ctx import cfg

_log = logging.getLogger("extract")

_IMAGE_SUFFIXES = (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff", ".gif")
# A page whose extracted text is shorter than this is treated as scanned (image
# only) and sent to OCR during ingestion.
_SCANNED_MIN_CHARS = 16


def _extract_xlsx(p: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(_sheet_text(ws))
    wb.close()
    return "\n\n".join(t for t in out if t.strip())


def _sheet_text(ws) -> str:
    rows: list[str] = []
    ncols = 0
    for row in ws.iter_rows(values_only=True):
        cells = [str(c) for c in row if c is not None]
        if cells:
            ncols = max(ncols, len(cells))
            rows.append("\t".join(cells))
    if not rows:
        return ""
    # State explicit counts up front so "how many rows/tracks" is answered from
    # metadata — the model can't reliably count hundreds of rows in-context (it
    # hallucinates and burns its token budget). Row 0 is (almost always) the header.
    nrows = len(rows)
    data_rows = max(nrows - 1, 0)
    stats = (
        f'[Sheet "{ws.title}": {nrows} rows ({data_rows} data rows after the header row), '
        f"{ncols} columns. Header row: {rows[0]}]"
    )
    return f"{stats}\n" + "\n".join(rows)


def _extract_csv(p: Path) -> str:
    """CSV as text, but prepended with a row/column-count header (same rationale as
    `_sheet_text`) so a counting question is answered from metadata, not by tallying
    lines. The body is the original text, unchanged."""
    import csv

    text = p.read_text(encoding="utf-8", errors="replace")
    if not text.strip():
        return text
    try:
        dialect = csv.Sniffer().sniff(text[:4096])
    except csv.Error:
        dialect = csv.excel
    rows = [r for r in csv.reader(text.splitlines(), dialect) if any(c.strip() for c in r)]
    if not rows:
        return text
    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    stats = (
        f"[CSV: {nrows} rows ({max(nrows - 1, 0)} data rows after the header row), "
        f'{ncols} columns. Header row: {chr(9).join(rows[0])}]'
    )
    return f"{stats}\n{text}"


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    parts.append(line)
    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    if notes:
        parts.append(f"[Notes] {notes}")
    return "\n".join(parts)


def _extract_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    return "\n\n".join(
        f"[Slide {i + 1}]\n{t}"
        for i, t in enumerate(_slide_text(s) for s in prs.slides)
        if t.strip()
    )


def _table_to_markdown(tbl: list) -> str:
    """Render one extracted table as a GitHub-style Markdown table. Row/column
    structure is exactly what pypdf flattens away (the s159A / table failure)."""
    rows = [
        [("" if c is None else str(c)).replace("\n", " ").strip() for c in row]
        for row in tbl
        if any(c is not None and str(c).strip() for c in row)
    ]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    out += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join(out)


def _pdf_pages_pdfplumber(p: Path) -> list[tuple[int, str]]:
    """Per-page PDF text via pdfplumber: layout-aware text for prose, plus every
    detected table appended as Markdown so row/column structure survives (pypdf
    flattens tables into a column-less list). Per-page → citation page_number kept."""
    import pdfplumber

    pages: list[tuple[int, str]] = []
    with pdfplumber.open(str(p)) as pdf:
        for i, page in enumerate(pdf.pages):
            try:
                text = page.extract_text() or ""
                tables = page.extract_tables() or []
            except Exception:
                _log.warning("pdfplumber failed on page %d of %s", i + 1, p.name)
                text, tables = "", []
            parts = [text] if text.strip() else []
            for tbl in tables:
                md = _table_to_markdown(tbl)
                if md:
                    parts.append("\n[Table]\n" + md)
            pages.append((i + 1, "\n\n".join(parts)))
    return pages


def _pdf_pages(p: Path) -> list[tuple[int, str]]:
    """Per-page PDF text. Table-aware pdfplumber when the `pdfplumber` knob is on
    (default); fast pypdf otherwise. Any pdfplumber error falls back to pypdf so a
    parser quirk never fails ingestion."""
    if cfg("pdfplumber", settings.ingest_pdfplumber):
        try:
            return _pdf_pages_pdfplumber(p)
        except Exception:
            _log.warning("pdfplumber unavailable on %s; falling back to pypdf", p.name)
    reader = pypdf.PdfReader(str(p))
    return [(i + 1, page.extract_text() or "") for i, page in enumerate(reader.pages)]


def extract(path: str, mime: str | None = None) -> str:
    """Synchronous native text extraction (no OCR). Use `extract_pages_ocr` for
    the ingest path that must OCR scanned PDFs / images."""
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix in (".txt", ".md", ".text") or (mime and mime.startswith("text/")):
        return p.read_text(encoding="utf-8", errors="replace")

    if suffix == ".docx":
        # Same flattener the tracked-change writer uses, so the assistant's view
        # matches the viewer's (tracked-changes flow §6 invariant).
        return tracked_changes.extract_body_text(str(p))

    if suffix == ".pdf":
        return "\n".join(t for _, t in _pdf_pages(p))

    if suffix in (".xlsx", ".xlsm"):
        return _extract_xlsx(p)

    if suffix == ".pptx":
        return _extract_pptx(p)

    if suffix == ".csv" or (mime and "csv" in mime.lower()):
        return _extract_csv(p)

    # Best-effort fallback: treat as text.
    return p.read_text(encoding="utf-8", errors="replace")


def extract_pages(path: str, mime: str | None = None) -> list[tuple[int, str]]:
    """Like `extract`, but preserving page boundaries for citation `page_number`.
    PDFs yield one (page_no, text) per page; XLSX one per sheet; PPTX one per
    slide; other formats are a single logical page. No OCR (sync)."""
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _pdf_pages(p)
    if suffix in (".xlsx", ".xlsm"):
        import openpyxl

        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        pages = [(i + 1, _sheet_text(ws)) for i, ws in enumerate(wb.worksheets)]
        wb.close()
        return pages or [(1, "")]
    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(p))
        return [(i + 1, _slide_text(s)) for i, s in enumerate(prs.slides)] or [(1, "")]
    return [(1, extract(path, mime))]


def _mime_for(path: Path, mime: str | None) -> str:
    if mime:
        return mime
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "application/pdf"
    return {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".webp": "image/webp", ".bmp": "image/bmp", ".tif": "image/tiff",
        ".tiff": "image/tiff", ".gif": "image/gif",
    }.get(suffix, "application/octet-stream")


async def extract_pages_ocr(path: str, mime: str | None = None) -> list[tuple[int, str]]:
    """The async ingest entry point. Native parsers for the office/text formats;
    OCR (via the OCR service) for images and scanned PDFs. A PDF with a real text
    layer is read natively; a PDF with (effectively) no text is sent to OCR whole
    — the service rasterises and segments it. Raises `ocr.OcrUnavailable` when a
    document needs OCR but it is off/unreachable, so the document is never indexed
    as empty (OCR-on-upload)."""
    from . import ocr

    p = Path(path)
    suffix = p.suffix.lower()

    # Loose image files → OCR (one logical page).
    if suffix in _IMAGE_SUFFIXES or (mime and mime.startswith("image/")):
        text = await ocr.ocr_bytes(p.read_bytes(), _mime_for(p, mime))
        return [(1, text)]

    # PDF: prefer the embedded text layer; fall back to OCR for scanned PDFs.
    if suffix == ".pdf":
        pages = extract_pages(path, mime)
        if sum(len(t.strip()) for _, t in pages) >= _SCANNED_MIN_CHARS:
            return pages  # has a usable text layer
        # Scanned (no text layer): hand the whole PDF to the OCR service.
        text = await ocr.ocr_bytes(p.read_bytes(), "application/pdf")
        return [(1, text)]

    # DOCX / XLSX / PPTX / plain text — native, no OCR.
    return extract_pages(path, mime)
