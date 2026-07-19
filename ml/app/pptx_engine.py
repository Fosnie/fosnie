# Copyright 2026 Private AI Ltd (SC881079)
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Presentation artefacts.
The model emits a **JSON slide spec** (taught by the `pptx-deck` skill); this builds
it into a real 16:9 `.pptx` via python-pptx — native editable text, native tables,
native OOXML charts and speaker notes.

The model chooses what each slide says and which archetype carries it; this module
owns all geometry. Autofit is deliberately not used: PowerPoint recomputes
`MSO_AUTO_SIZE` when the file is opened, so a size written here would be a guess the
client silently overrides. Overflow is handled instead by character budgets, a font
ladder, and hard splits onto `(continued)` slides.

A drafter that emits Markdown rather than JSON is also accepted (`_parse`), so the
fallback document path still produces a deck rather than an error."""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path

from .config import settings

log = logging.getLogger(__name__)

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# --- design system -----------------------------------------------------------

# Deck geometry: 16:9 widescreen with a half-inch margin on every side.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.5
CONTENT_W_IN = SLIDE_W_IN - 2 * MARGIN_IN

# Neutral defaults, matched to the HTML artefact theme so the formats read as one
# family. A deployment overrides the two brand colours and the font through config.
# A deck names exactly one typeface (no CSS-style fallback list), so the default is
# the second entry of the theme's stack rather than the first: it is present on
# Windows and macOS, and metric-compatible on Linux. Naming a face the recipient
# lacks means silent substitution and a deck that reflows on their machine.
DEFAULT_FONT = "Arial"
INK = "1F2933"
MUTED = "8A97A4"
LINE = "E2E6EA"
ZEBRA = "F0F2F5"  # the subtle raised surface, used for even table rows
PRIMARY = "2F6DB5"
ACCENT = "1F9D8F"
# Categorical chart palette (the same six the HTML theme uses).
SERIES_COLOURS = ["2F6DB5", "1F9D8F", "C98A1A", "7A5BB0", "C0492F", "4A8DB2"]

LAYOUTS = frozenset({
    "title", "section", "bullets", "two_column", "stat", "table", "chart", "quote", "closing",
})

CHART_TYPES = frozenset({"bar", "column", "line", "pie", "doughnut"})

# Soft budgets: exceeding one steps the font down rather than failing.
TITLE_BUDGET = 60
HEADING_BUDGET = 70
BULLET_BUDGET = 90
# Hard splits: beyond these a slide becomes two, because no font size rescues them.
BULLETS_PER_SLIDE = 5
BULLETS_SPLIT_AT = 6
ROWS_PER_SLIDE = 8
ROWS_SPLIT_AT = 8
MAX_CATEGORIES = 8
MAX_SERIES = 4
# Last-resort clip so a runaway field cannot run off the slide entirely.
HARD_CLIP = 400


def _hex(configured: str, default: str) -> str:
    """A six-digit hex colour from config, or the default. A malformed value is
    logged and ignored rather than failing the generation — a branding typo must not
    cost the user their deck."""
    v = (configured or "").strip().lstrip("#").upper()
    if not v:
        return default
    if re.fullmatch(r"[0-9A-F]{6}", v):
        return v
    log.warning("ignoring malformed pptx brand colour %r", configured)
    return default


def _clip(text, budget: int = HARD_CLIP) -> str:
    s = " ".join(str(text if text is not None else "").split())
    return s if len(s) <= budget else s[: budget - 1].rstrip() + "…"


def _ladder(length: int, steps: list[tuple[int, int]], smallest: int) -> int:
    """Pick a point size: the first step whose character budget still holds."""
    for budget, size in steps:
        if length <= budget:
            return size
    return smallest


# --- parsing -----------------------------------------------------------------


def _parse(content: str, title: str) -> dict:
    """Accept `{slides:[…]}`, a bare list of slides, or a single slide object. When
    the content is not JSON at all, read it as Markdown (the drafter path). Raises
    ValueError (→ HTTP 400) on anything that yields no slides."""
    raw = (content or "").strip()
    if not raw:
        raise ValueError("pptx content is empty")

    try:
        spec = json.loads(raw)
    except Exception:
        slides = _from_markdown(raw, title)
        if not slides:
            raise ValueError(
                "pptx content must be a JSON slide spec (or Markdown with headings)"
            ) from None
        return {"slides": slides}

    if isinstance(spec, list):
        spec = {"slides": spec}
    if not isinstance(spec, dict):
        raise ValueError("pptx spec must be a JSON object or array")
    if "layout" in spec:  # a single slide sent on its own
        spec = {"slides": [spec]}
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("pptx spec needs a non-empty 'slides' array")
    for s in slides:
        if not isinstance(s, dict):
            raise ValueError("each slide must be an object")
        layout = s.get("layout", "bullets")
        if layout not in LAYOUTS:
            raise ValueError(f"unknown slide layout: {layout}")
    return {"slides": slides}


_INLINE = re.compile(r"\*\*|__|`|\*(?!\s)|(?<!\s)\*")
_LINK = re.compile(r"\[([^\]]*)\]\([^)]*\)")


def _plain(text: str) -> str:
    """Strip the inline Markdown emphasis that would otherwise render literally."""
    return _INLINE.sub("", _LINK.sub(r"\1", text)).strip()


_BULLET = re.compile(r"^\s*(?:[-*+]|\d+[.)])\s+(.*)$")


def _from_markdown(md: str, title: str) -> list[dict]:
    """Read Markdown as a deck: the H1 (or the artefact title) opens it, each H2 or
    thematic break starts a slide, and list items and paragraphs become bullets.
    Robustness for the fallback path, not a substitute for the slide spec."""
    slides: list[dict] = []
    heading = ""
    bullets: list[str] = []
    deck_title = ""
    subtitle = ""

    def flush() -> None:
        nonlocal heading, bullets
        if heading or bullets:
            slides.append({
                "layout": "bullets",
                "title": heading or (deck_title or title or "Overview"),
                "bullets": bullets[:],
            })
        heading, bullets = "", []

    for line in md.replace("\r\n", "\n").split("\n"):
        s = line.strip()
        if not s:
            continue
        if s.startswith("# ") and not deck_title:
            deck_title = _plain(s[2:])
            continue
        if s.startswith("## ") or (s.startswith("#") and not s.startswith("###")):
            flush()
            heading = _plain(s.lstrip("#").strip())
            continue
        if s.startswith("### "):
            bullets.append(_plain(s[4:]))
            continue
        if set(s) <= {"-", "*", "_"} and len(s) >= 3:  # thematic break
            flush()
            continue
        m = _BULLET.match(s)
        text = _plain(m.group(1)) if m else _plain(s)
        if not text:
            continue
        if not slides and not heading and not bullets and deck_title and not subtitle:
            subtitle = text  # the lead paragraph under the H1
            continue
        bullets.append(text)
    flush()

    if deck_title or title:
        opener: dict = {"layout": "title", "title": deck_title or title}
        if subtitle:
            opener["subtitle"] = subtitle
        slides.insert(0, opener)
    return slides


def _expand(slides: list[dict]) -> list[dict]:
    """Split what no font size can rescue: long bullet lists and long tables become
    a run of `(continued)` slides. Speaker notes stay with the first part."""
    out: list[dict] = []
    for s in slides:
        layout = s.get("layout", "bullets")
        if layout in ("bullets", "closing"):
            items = s.get("bullets") or []
            if isinstance(items, list) and len(items) > BULLETS_SPLIT_AT:
                for i in range(0, len(items), BULLETS_PER_SLIDE):
                    part = dict(s)
                    part["bullets"] = items[i : i + BULLETS_PER_SLIDE]
                    if i:
                        part["title"] = f"{s.get('title', '')} (continued)"
                        part.pop("notes", None)
                    out.append(part)
                continue
        elif layout == "table":
            rows = s.get("rows") or []
            if isinstance(rows, list) and len(rows) > ROWS_SPLIT_AT:
                for i in range(0, len(rows), ROWS_PER_SLIDE):
                    part = dict(s)
                    part["rows"] = rows[i : i + ROWS_PER_SLIDE]
                    if i:
                        part["title"] = f"{s.get('title', '')} (continued)"
                        part.pop("notes", None)
                        part.pop("caption", None)
                    out.append(part)
                continue
        out.append(s)
    return out


# --- rendering ---------------------------------------------------------------


def build(content: str, title: str, out_path: str) -> dict:
    spec = _parse(content, title)

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    font = (settings.pptx_font or "").strip() or DEFAULT_FONT
    primary = RGBColor.from_string(_hex(settings.pptx_colour_primary, PRIMARY))
    accent = RGBColor.from_string(_hex(settings.pptx_colour_accent, ACCENT))
    theme = {
        "font": font,
        "primary": primary,
        "accent": accent,
        "ink": RGBColor.from_string(INK),
        "muted": RGBColor.from_string(MUTED),
        "line": RGBColor.from_string(LINE),
        "white": RGBColor.from_string("FFFFFF"),
        "series": [RGBColor.from_string(c) for c in SERIES_COLOURS],
    }

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    section_no = 0
    slides = _expand(spec["slides"])
    for i, s in enumerate(slides):
        layout = s.get("layout", "bullets")
        slide = prs.slides.add_slide(blank)
        if layout == "section":
            section_no += 1
        _render(slide, layout, s, theme, section_no)

        if layout != "title":
            _footer(slide, i + 1, theme)
        notes = s.get("notes")
        if isinstance(notes, str) and notes.strip():
            slide.notes_slide.notes_text_frame.text = notes.strip()

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {"path": str(out), "mime": PPTX_MIME}


def _render(slide, layout: str, s: dict, theme: dict, section_no: int) -> None:
    if layout == "title":
        _title_slide(slide, s, theme)
    elif layout == "section":
        _section_slide(slide, s, theme, section_no)
    elif layout == "quote":
        _quote_slide(slide, s, theme)
    elif layout == "two_column":
        _two_column_slide(slide, s, theme)
    elif layout == "stat":
        _stat_slide(slide, s, theme)
    elif layout == "table":
        _table_slide(slide, s, theme)
    elif layout == "chart":
        _chart_slide(slide, s, theme)
    else:  # bullets, closing
        _bullets_slide(slide, s, theme)


def _textbox(slide, left, top, width, height):
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _write(tf, text: str, theme: dict, size: int, *, bold=False, colour=None,
           align=None, para=None, space_after=0, italic=False):
    """Write one line into a text frame as an explicitly styled run. Every run
    carries its own font because a blank layout supplies no placeholder styling."""
    from pptx.util import Pt

    p = para if para is not None else tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = theme["font"]
    run.font.color.rgb = colour if colour is not None else theme["ink"]
    if align is not None:
        p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    return p


def _bar(slide, left, top, width, theme, colour_key="accent"):
    """The short rule that anchors a centred slide."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme[colour_key]
    shape.line.fill.background()
    shape.shadow.inherit = False


def _footer(slide, number: int, theme: dict) -> None:
    from pptx.enum.text import PP_ALIGN

    tf = _textbox(slide, SLIDE_W_IN - 1.6, SLIDE_H_IN - 0.62, 1.1, 0.35)
    _write(tf, str(number), theme, 10, colour=theme["muted"], align=PP_ALIGN.RIGHT)


def _heading(slide, s: dict, theme: dict) -> None:
    """The left-aligned slide heading. It states the idea; the body evidences it."""
    text = _clip(s.get("title", ""))
    if not text:
        return
    size = _ladder(len(text), [(HEADING_BUDGET, 28), (110, 24)], 22)
    tf = _textbox(slide, MARGIN_IN, 0.55, CONTENT_W_IN, 1.0)
    _write(tf, text, theme, size, bold=True)
    _bar(slide, MARGIN_IN, 1.52, 1.4, theme)


def _title_slide(slide, s: dict, theme: dict) -> None:
    from pptx.enum.text import PP_ALIGN

    title = _clip(s.get("title", ""))
    _bar(slide, (SLIDE_W_IN - 1.6) / 2, 2.35, 1.6, theme)
    tf = _textbox(slide, 1.2, 2.65, SLIDE_W_IN - 2.4, 1.5)
    _write(tf, title, theme, _ladder(len(title), [(TITLE_BUDGET, 40)], 32), bold=True,
           align=PP_ALIGN.CENTER)
    subtitle = _clip(s.get("subtitle", ""))
    if subtitle:
        tf2 = _textbox(slide, 1.2, 4.25, SLIDE_W_IN - 2.4, 0.9)
        _write(tf2, subtitle, theme, 20, colour=theme["muted"], align=PP_ALIGN.CENTER)


def _section_slide(slide, s: dict, theme: dict, section_no: int) -> None:
    from pptx.enum.text import PP_ALIGN

    tf = _textbox(slide, 1.2, 2.5, SLIDE_W_IN - 2.4, 0.9)
    _write(tf, f"{section_no:02d}", theme, 32, bold=True, colour=theme["accent"],
           align=PP_ALIGN.CENTER)
    title = _clip(s.get("title", ""))
    tf2 = _textbox(slide, 1.2, 3.35, SLIDE_W_IN - 2.4, 1.4)
    _write(tf2, title, theme, _ladder(len(title), [(50, 36)], 28), bold=True,
           align=PP_ALIGN.CENTER)


def _quote_slide(slide, s: dict, theme: dict) -> None:
    text = _clip(s.get("text", ""))
    _bar(slide, 1.2, 2.15, 1.2, theme)
    tf = _textbox(slide, 1.2, 2.5, SLIDE_W_IN - 2.4, 2.6)
    _write(tf, text, theme, _ladder(len(text), [(120, 30), (220, 26)], 22), italic=True)
    attribution = _clip(s.get("attribution", ""))
    if attribution:
        tf2 = _textbox(slide, 1.2, 5.2, SLIDE_W_IN - 2.4, 0.6)
        _write(tf2, attribution, theme, 16, colour=theme["muted"])


def _bullet_items(raw) -> list[tuple[str, list[str]]]:
    """Normalise bullets to (text, subs). An item is a string or {text, sub}."""
    items: list[tuple[str, list[str]]] = []
    for b in raw if isinstance(raw, list) else []:
        if isinstance(b, dict):
            subs = b.get("sub") or []
            items.append((
                _clip(b.get("text", "")),
                [_clip(x) for x in subs if str(x).strip()] if isinstance(subs, list) else [],
            ))
        elif str(b).strip():
            items.append((_clip(b), []))
    return items


def _bullet_glyph(para, theme: dict, char: str, level: int) -> None:
    """Turn a paragraph into a real bulleted list item. DrawingML carries the bullet
    as paragraph formatting, so the glyph is not part of the text: wrapped lines hang
    correctly and the recipient gets a genuine list to edit in PowerPoint. python-pptx
    has no high-level API for this, hence the direct pPr. Child order matters to the
    schema, so these are appended after any spacing python-pptx has already written."""
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches

    hang = Inches(0.28)
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", str(Emu(hang * (level + 1))))
    pPr.set("indent", str(Emu(-hang)))
    font = pPr.makeelement(qn("a:buFont"), {"typeface": theme["font"]})
    glyph = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(font)
    pPr.append(glyph)


def _bullets_into(tf, items, theme: dict, size: int) -> None:
    from pptx.util import Pt

    first = True
    for text, subs in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _write(tf, text, theme, size, para=p, space_after=10)
        _bullet_glyph(p, theme, "•", 0)
        for sub in subs[:3]:
            sp = tf.add_paragraph()
            sp.level = 1
            _write(tf, sub, theme, max(size - 2, 12), para=sp,
                   colour=theme["muted"], space_after=6)
            sp.space_before = Pt(2)
            _bullet_glyph(sp, theme, "–", 1)


def _bullets_slide(slide, s: dict, theme: dict) -> None:
    _heading(slide, s, theme)
    items = _bullet_items(s.get("bullets"))
    if not items:
        return
    longest = max(len(t) for t, _ in items)
    crowded = len(items) > 5 or longest > BULLET_BUDGET
    size = 18 if not crowded else (16 if longest <= 140 and len(items) <= 6 else 14)
    tf = _textbox(slide, MARGIN_IN, 1.85, CONTENT_W_IN, 5.0)
    _bullets_into(tf, items, theme, size)


def _two_column_slide(slide, s: dict, theme: dict) -> None:
    _heading(slide, s, theme)
    gap = 0.4
    col_w = (CONTENT_W_IN - gap) / 2
    for idx, key in enumerate(("left", "right")):
        col = s.get(key) or {}
        if not isinstance(col, dict):
            continue
        left = MARGIN_IN + idx * (col_w + gap)
        top = 1.85
        heading = _clip(col.get("heading", ""))
        if heading:
            tf = _textbox(slide, left, top, col_w, 0.5)
            _write(tf, heading, theme, 18, bold=True, colour=theme["primary"])
            top += 0.6
        items = _bullet_items(col.get("bullets"))
        if not items:
            continue
        longest = max(len(t) for t, _ in items)
        size = 16 if longest <= 90 and len(items) <= 4 else 14
        _bullets_into(_textbox(slide, left, top, col_w, 6.5 - top), items, theme, size)


def _stat_slide(slide, s: dict, theme: dict) -> None:
    from pptx.enum.text import PP_ALIGN

    _heading(slide, s, theme)
    stats = [x for x in (s.get("stats") or []) if isinstance(x, dict)][:4]
    if not stats:
        return
    gap = 0.4
    col_w = (CONTENT_W_IN - gap * (len(stats) - 1)) / len(stats)
    for i, stat in enumerate(stats):
        left = MARGIN_IN + i * (col_w + gap)
        value = _clip(stat.get("value", ""), 40)
        tf = _textbox(slide, left, 2.5, col_w, 1.5)
        _write(tf, value, theme, _ladder(len(value), [(6, 60), (12, 44)], 32), bold=True,
               colour=theme["accent"], align=PP_ALIGN.CENTER)
        label = _clip(stat.get("label", ""))
        tf2 = _textbox(slide, left, 4.1, col_w, 1.2)
        _write(tf2, label, theme, 14, colour=theme["muted"], align=PP_ALIGN.CENTER)


def _caption(slide, s: dict, theme: dict, top: float) -> None:
    text = _clip(s.get("caption", ""))
    if text:
        _write(_textbox(slide, MARGIN_IN, top, CONTENT_W_IN, 0.5), text, theme, 11,
               colour=theme["muted"])


def _table_slide(slide, s: dict, theme: dict) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    _heading(slide, s, theme)
    columns = [_clip(c, 80) for c in (s.get("columns") or [])]
    rows = [r for r in (s.get("rows") or []) if isinstance(r, list)]
    if not columns and not rows:
        raise ValueError("a table slide needs 'columns' or 'rows'")
    width = len(columns) or max((len(r) for r in rows), default=1)
    if not columns:
        columns = [""] * width

    top = 1.9
    row_h = 0.42
    shape = slide.shapes.add_table(
        len(rows) + 1, width, Inches(MARGIN_IN), Inches(top),
        Inches(CONTENT_W_IN), Inches(row_h * (len(rows) + 1)),
    )
    table = shape.table
    table.first_row = True
    for c in range(width):
        table.columns[c].width = Inches(CONTENT_W_IN / width)

    longest = max([len(x) for x in columns] + [len(_clip(v, 80)) for r in rows for v in r] + [1])
    size = 14 if width <= 4 and longest <= 28 else 12
    zebra = RGBColor.from_string(ZEBRA)

    for c, header in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["primary"]
        cell.text_frame.word_wrap = True
        _write(cell.text_frame, header, theme, size, bold=True, colour=theme["white"])
    for r, row in enumerate(rows, start=1):
        table.rows[r].height = Inches(row_h)
        for c in range(width):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["white"] if r % 2 else zebra
            cell.text_frame.word_wrap = True
            value = _clip(row[c], 80) if c < len(row) else ""
            _write(cell.text_frame, value, theme, size)

    _caption(slide, s, theme, top + row_h * (len(rows) + 1) + 0.15)


def _chart_slide(slide, s: dict, theme: dict) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    _heading(slide, s, theme)
    spec = s.get("chart")
    if not isinstance(spec, dict):
        raise ValueError("a chart slide needs a 'chart' object")
    ctype = str(spec.get("type", "column")).lower()
    if ctype not in CHART_TYPES:
        raise ValueError(f"unknown chart type: {ctype} (bar|column|line|pie|doughnut)")

    categories = [_clip(c, 40) for c in (spec.get("categories") or [])]
    if not categories:
        raise ValueError("a chart needs 'categories'")
    if len(categories) > MAX_CATEGORIES:
        raise ValueError(f"a chart takes at most {MAX_CATEGORIES} categories")

    series = [x for x in (spec.get("series") or []) if isinstance(x, dict)]
    if not series:
        raise ValueError("a chart needs at least one series")
    if len(series) > MAX_SERIES:
        raise ValueError(f"a chart takes at most {MAX_SERIES} series")
    if ctype in ("pie", "doughnut") and len(series) != 1:
        raise ValueError(f"a {ctype} chart takes exactly one series")

    data = CategoryChartData()
    data.categories = categories
    for one in series:
        values = one.get("values") or []
        if not isinstance(values, list):
            raise ValueError("chart series 'values' must be an array of numbers")
        nums = []
        for v in values[: len(categories)]:
            try:
                nums.append(float(v))
            except (TypeError, ValueError):
                raise ValueError(f"chart values must be numbers, got {v!r}") from None
        nums += [None] * (len(categories) - len(nums))
        data.add_series(_clip(one.get("name", "Series"), 60), nums)

    kinds = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    has_caption = bool(str(s.get("caption", "")).strip())
    height = 4.55 if has_caption else 4.9
    chart = slide.shapes.add_chart(
        kinds[ctype], Inches(MARGIN_IN), Inches(1.85), Inches(CONTENT_W_IN), Inches(height), data
    ).chart

    chart.has_title = False
    chart.font.size = Pt(12)
    chart.font.name = theme["font"]
    chart.font.color.rgb = theme["ink"]
    _style_series(chart, ctype, theme)

    if ctype in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.category_axis.has_major_gridlines = False
        chart.value_axis.has_major_gridlines = True

    _caption(slide, s, theme, 1.85 + height + 0.1)


def _style_series(chart, ctype: str, theme: dict) -> None:
    """Paint the series from the deck palette. Pie and doughnut colour their points
    rather than the series, since one series IS the whole chart."""
    colours = theme["series"]
    if ctype in ("pie", "doughnut"):
        plot = chart.plots[0]
        for i, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colours[i % len(colours)]
        return
    for i, one in enumerate(chart.series):
        colour = colours[i % len(colours)]
        if ctype == "line":
            one.format.line.color.rgb = colour
            one.smooth = False
        else:
            one.format.fill.solid()
            one.format.fill.fore_color.rgb = colour
